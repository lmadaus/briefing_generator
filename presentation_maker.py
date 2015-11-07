#!/usr/bin/env python
from __future__ import print_function, division
from pptx import Presentation
from pptx.util import Inches, Pt
#from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from datetime import time as dtime
from datetime import datetime, timedelta
import os, re
import urllib

# Figure out the current time
now = datetime.now()
utcnow = datetime.utcnow()
nowdate = now.date()
utcnowdate = utcnow.date()
# Now, based on the current time, decide which presentation we are doing
# (18Z morning or 4Z evening) and which should be the most recent available
# model run
if utcnow.hour < 16 and utcnow.hour >= 4:
    model_init_date = datetime.combine(utcnowdate, dtime(0,0))
    present_date = datetime.combine(utcnowdate,dtime(5,0))
elif utcnow.hour < 4:

    model_init_date = datetime.combine(utcnowdate-timedelta(hours=24),dtime(12,0))
    present_date = datetime.combine(nowdate,dtime(18,0))
else:
    model_init_date = datetime.combine(utcnowdate,dtime(12,0))
    present_date = datetime.combine(nowdate,dtime(18,0))    

print("Model init date:", model_init_date)
print("Presentation date:", present_date)
# Uncomment the following line to manually specify the
# model data to grab
#model_init_date=datetime(2015,11,5,12)


# Set the directory where the output files will be put
presentation_path = '.'
# Set the path to the model images
model_path='http://www.atmos.washington.edu/~lmadaus/olympex/index.php?init={:%Y%m%d%H}&product={:s}&start={:d}'

# Different default possibilities for the slide layout
layout = {'Title Slide' : 0,
          'Bullet Slide' : 1,
          'Segue' : 2,
          'Side By Side' : 3,
          'Title Alone' : 5,
          'Blank Slide' : 6,
          'Picture with Caption' : 8,
          
         }

# The pattern for the image paths is:
#   'Product name' : (path to file or web address, filename ending )
img_paths = {
             #'IR+500mb' : ('http://www.atmos.washington.edu/images/sat_upr/YYYYMMDDHH00_500mb.gif','500mb.gif',),
             #'Water Vapor' : ('http://www.atmos.washington.edu/images/sat_common/YYYYMMDDHH00_wv.gif','wv.gif'),
             'IR+500mb' : ('http://www.atmos.washington.edu/cgi-bin/latest.cgi?sat_500+-notitle','500mb.gif',),
             'Water Vapor' : ('http://www.atmos.washington.edu/cgi-bin/latest.cgi?wv_common_full+-notitle','wv.gif'),
             'OPC Surface Analys.' : ('http://www.opc.ncep.noaa.gov/P_e_sfc_color.png',None),
             'WRF 500mb Day 0' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxLG_500mb.YYYYMMDDHH.f009.png', 'wrf_500mb_day0.png'),
             'WRF 500mb Day 1' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxLG_500mb.YYYYMMDDHH.f024.png', 'wrf_500mb_day1.png'),
             'WRF 500mb Day 2' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxLG_500mb.YYYYMMDDHH.f048.png', 'wrf_500mb_day2.png'),
             'WRF SLP Day 0' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxLG_surface.YYYYMMDDHH.f009.png', 'wrf_sfc_day0.png'),
             'WRF SLP Day 1' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxLG_surface.YYYYMMDDHH.f024.png', 'wrf_sfc_day1.png'),
             'WRF SLP Day 2' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxLG_surface.YYYYMMDDHH.f048.png', 'wrf_sfc_day2.png'),

             'WRF Melt. Level Day 0' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxSM_melt_level.YYYYMMDDHH.f009.png', 'wrf_melt_level_day0.png'),
             'WRF Melt. Level Day 1' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxSM_melt_level.YYYYMMDDHH.f024.png', 'wrf_melt_level_day1.png'),     
             'WRF Melt. Level Day 2' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxSM_melt_level.YYYYMMDDHH.f048.png', 'wrf_melt_level_day2.png'),     

             'WRF 12hr Prcp Day 0' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxLG_precip12hr.YYYYMMDDHH.f012.png', 'wrf_precip_large_day0.png'),
             'WRF 12hr Prcp (4km) Day 0' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxSM_precip12hr.YYYYMMDDHH.f012.png', 'wrf_precip_small_day0.png'),
             'WRF 12hr Prcp Day 1' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxLG_precip12hr.YYYYMMDDHH.f024.png', 'wrf_precip_large_day1.png'),
             'WRF 12hr Prcp (4km) Day 1' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxSM_precip12hr.YYYYMMDDHH.f024.png', 'wrf_precip_small_day1.png'),
             'WRF 12hr Prcp Day 2' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxLG_precip12hr.YYYYMMDDHH.f048.png', 'wrf_precip_large_day2.png'),
             'WRF 12hr Prcp (4km) Day 2' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxSM_precip12hr.YYYYMMDDHH.f048.png', 'wrf_precip_small_day2.png'),


             'WRF 3hr Prcp Day 0' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxLG_precip3hr.YYYYMMDDHH.f009.png', 'wrf_precip03_large_day0.png'),
             'WRF 3hr Prcp (4km) Day 0' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxSM_precip3hr.YYYYMMDDHH.f009.png', 'wrf_precip03_small_day0.png'),
             'WRF 3hr Prcp Day 1' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxLG_precip3hr.YYYYMMDDHH.f024.png', 'wrf_precip03_large_day1.png'),
             'WRF 3hr Prcp (4km) Day 1' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxSM_precip3hr.YYYYMMDDHH.f024.png', 'wrf_precip03_small_day1.png'),
             'WRF 3hr Prcp (4km) Day 1a' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxSM_precip3hr.YYYYMMDDHH.f018.png', 'wrf_precip03_small_day1a.png'),
             'WRF 3hr Prcp (4km) Day 1b' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxSM_precip3hr.YYYYMMDDHH.f024.png', 'wrf_precip03_small_day1b.png'),
             'WRF 3hr Prcp (4km) Day 1c' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxSM_precip3hr.YYYYMMDDHH.f030.png', 'wrf_precip03_small_day1c.png'),
             'WRF 3hr Prcp (4km) Day 1d' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxSM_precip3hr.YYYYMMDDHH.f036.png', 'wrf_precip03_small_day1d.png'),

             'WRF 3hr Prcp Day 2' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxLG_precip3hr.YYYYMMDDHH.f048.png', 'wrf_precip03_large_day2.png'),
             'WRF 3hr Prcp (4km) Day 2' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxSM_precip3hr.YYYYMMDDHH.f048.png', 'wrf_precip03_small_day2.png'),
             'WRF 3hr Prcp (4km) Day 2a' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxSM_precip3hr.YYYYMMDDHH.f042.png', 'wrf_precip03_small_day2a.png'),
             'WRF 3hr Prcp (4km) Day 2b' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxSM_precip3hr.YYYYMMDDHH.f048.png', 'wrf_precip03_small_day2b.png'),
             'WRF 3hr Prcp (4km) Day 2c' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxSM_precip3hr.YYYYMMDDHH.f054.png', 'wrf_precip03_small_day2c.png'),
             'WRF 3hr Prcp (4km) Day 2d' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxSM_precip3hr.YYYYMMDDHH.f060.png', 'wrf_precip03_small_day2d.png'),

             'WRF 3hr Prcp (4km) Day 3a' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxSM_precip3hr.YYYYMMDDHH.f066.png', 'wrf_precip03_small_day3a.png'),
             'WRF 3hr Prcp (4km) Day 3b' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxSM_precip3hr.YYYYMMDDHH.f072.png', 'wrf_precip03_small_day3b.png'),
             'WRF 3hr Prcp (4km) Day 3c' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxSM_precip3hr.YYYYMMDDHH.f078.png', 'wrf_precip03_small_day3c.png'),
             'WRF 3hr Prcp (4km) Day 3d' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxSM_precip3hr.YYYYMMDDHH.f084.png', 'wrf_precip03_small_day3d.png'),

             'WRF 10m Wind Day 0' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxLG_wssfc.YYYYMMDDHH.f009.png', 'wrf_wssfc_large_day0.png'),
             'WRF 10m Wind (4km) Day 0' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxSM_wssfc.YYYYMMDDHH.f009.png', 'wrf_wssfc_small_day0.png'),
             'WRF 10m Wind Day 1' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxLG_wssfc.YYYYMMDDHH.f024.png', 'wrf_wssfc_large_day1.png'),
             'WRF 10m Wind (4km) Day 1' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxSM_wssfc.YYYYMMDDHH.f024.png', 'wrf_wssfc_small_day1.png'),
             'WRF 10m Wind Day 2' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxLG_wssfc.YYYYMMDDHH.f048.png', 'wrf_wssfc_large_day2.png'),
             'WRF 10m Wind (4km) Day 2' : ('http://www.atmos.washington.edu/~lmadaus/olympex/wrf_plots/YYYYMMDDHH/opxSM_wssfc.YYYYMMDDHH.f048.png', 'wrf_wssfc_small_day2.png'),             
             
             'NWS PacNW Radar' : ('http://www.atmos.washington.edu/~lmadaus/olympex/radar/YYYYMMDDHH00_bref.png','radar.png'),    
             'KUIL Latest Sound.' : ('http://www.atmos.washington.edu/~lmadaus/olympex/soundings/KUIL_YYYYMMDDHH_snd.png','sounding.png'),
             'GFS 500mb Day 2' : ('','gfs_500_day3.gif'),
             'NAEFS SLP and Spread Day 2' : ('http://collaboration.cmc.ec.gc.ca/cmc/ensemble/cartes/data/cartes/PN/CMC_NCEP','naefs_slp_day2.gif'),
             'NAEFS SLP and Spread Day 3' : ('http://collaboration.cmc.ec.gc.ca/cmc/ensemble/cartes/data/cartes/PN/CMC_NCEP','naefs_slp_day3.gif'),
             'NAEFS SLP and Spread Day 4' : ('http://collaboration.cmc.ec.gc.ca/cmc/ensemble/cartes/data/cartes/PN/CMC_NCEP','naefs_slp_day4.gif'),
             'NAEFS SLP and Spread Day 5' : ('http://collaboration.cmc.ec.gc.ca/cmc/ensemble/cartes/data/cartes/PN/CMC_NCEP','naefs_slp_day5.gif'),

             'NAEFS 500mb and Spread Day 2' : ('http://collaboration.cmc.ec.gc.ca/cmc/ensemble/cartes/data/cartes/GZ500/CMC_NCEP','naefs_500mb_day2.gif'),
             'NAEFS 500mb and Spread Day 3' : ('http://collaboration.cmc.ec.gc.ca/cmc/ensemble/cartes/data/cartes/GZ500/CMC_NCEP','naefs_500mb_day3.gif'),
             'NAEFS 500mb and Spread Day 4' : ('http://collaboration.cmc.ec.gc.ca/cmc/ensemble/cartes/data/cartes/GZ500/CMC_NCEP','naefs_500mb_day4.gif'),
             'NAEFS 500mb and Spread Day 5' : ('http://collaboration.cmc.ec.gc.ca/cmc/ensemble/cartes/data/cartes/GZ500/CMC_NCEP','naefs_500mb_day5.gif'),

            }

def build_presentation(model_init_date, present_date):
    """
    Function to build the final presentation
    Inputs:
    model_init_date --> the initialization date for the model graphics to grab
    present_date --> The actual key time of the presentation
    """
    # Switch to the presentation directory
    basedir = os.getcwd()
    os.chdir(presentation_path)
   
    # The actual presentation time will be 6 hours later
    #present_date = model_init_date + timedelta(hours=6)

    # Make a new subdirectory for this presentation
    if os.path.exists('./{:%Y%m%d%H}'.format(present_date)):
        os.system('rm ./{:%Y%m%d%H}/*'.format(present_date))
    else:     
        os.mkdir('./{:%Y%m%d%H}'.format(present_date))
    os.chdir('./{:%Y%m%d%H}'.format(present_date))

    # Make a new presentation
    #try:
    #prs = Presentation(os.path.join(basedir,'default.pptx'))
    #except:
    prs = Presentation()
    
    # Make the title slide
    # Choose a title slide layout
    title_slide_layout = prs.slide_layouts[layout['Title Slide']]
    # Add the slide to the presentation
    title_slide = prs.slides.add_slide(title_slide_layout)
    # Change the title
    if model_init_date.hour < 12:
        title_slide.shapes.title.text = "Evening Weather Update"

    else:  
        title_slide.shapes.title.text = "Morning Weather Briefing"
    # Subtitle is a "placeholder" object
    title_slide.placeholders[1].text = "{:%d %b %Y / %H00Z}\nForecaster Name".format(present_date)

    
    # Make Prev weather bumper
    prs = bumper_slide(prs, 'Past 24 hours', present_date-timedelta(days=1))
    # Blank slide here
    prs = full_summary(prs, 'Summary of Prev. 24 Hours')



    # Make current weather bumper
    prs = bumper_slide(prs, 'Current Weather', present_date)
    # Call the above function to add full-slide images
    # Here for IR+500
    prs = full_slide_image(prs, 'IR+500mb', model_init_date, width=9, link='http://www.atmos.washington.edu/~ovens/wxloop.cgi?sat_500+/2d/')
    # And for Water Vapor
    prs = full_slide_image(prs, 'Water Vapor', model_init_date, width=9, link='http://www.atmos.washington.edu/~ovens/wxloop.cgi?wv_common+/48h/')
    # Here for the web-based OPC surface analysis
    prs = full_slide_image(prs, 'OPC Surface Analys.', model_init_date, link=False)

    # Model verification?    

    
    # Regional Radar
    prs = full_slide_image(prs, 'NWS PacNW Radar', model_init_date, link='http://www.atmos.washington.edu/~lmadaus/olympex/index.php?product=radar')
    # Current Sounding
    prs = full_slide_image(prs, 'KUIL Latest Sound.', model_init_date)
    
    # Current airport conditions--> vis, wind dir, wind speed (and at NPOL)
    #prs = full_summary(prs, 'Current Airport Conditions')
    prs = wxdata_slide(prs, 'Current Airport Conditions', type='METAR', locs=['KPAE','KTCM','KHQM'])
    
    # Bumper into next 24 hour forecast
    prs = bumper_slide(prs, 'Forecast: Day 0', model_init_date)
    # Sort out the days now
    
    day0_start = datetime.combine(present_date,dtime(0,0))
    day1_start = day0_start + timedelta(hours=24)
    day2_start = day0_start + timedelta(hours=48)
    day3_start = day0_start + timedelta(hours=72)
    day0_ftime = day0_start + timedelta(hours=21)
    day1_ftime = day1_start + timedelta(hours=12)
    day2_ftime = day2_start + timedelta(hours=12)
    day3_ftime = day3_start + timedelta(hours=12)
    #day1_ftime = day1_ftime.replace(hour=0, minute=0, second=0)
    
   
    # WRF Image --> 500mb Vort
    prs = full_slide_image(prs, 'WRF 500mb Day 0', model_init_date, day0_ftime, link=model_path.format(model_init_date,'opxLG_500mb',1))
    
    # WRF SLP
    prs = full_slide_image(prs, 'WRF SLP Day 0', model_init_date, day0_ftime, link=model_path.format(model_init_date,'opxLG_surface',1))
    
    # WRF Melting level
    prs = full_slide_image(prs, 'WRF Melt. Level Day 0', model_init_date, day0_ftime, link=model_path.format(model_init_date,'opxSM_melt_level',1))
       
    # WRF zoom Precip
    prs = full_slide_image(prs, 'WRF 3hr Prcp (4km) Day 0', model_init_date, day0_ftime, link=model_path.format(model_init_date,'opxSM_precip3hr',1)) 
      
    # WRF zoom 10m Winds
    prs = full_slide_image(prs, 'WRF 10m Wind (4km) Day 0', model_init_date, day0_ftime, link=model_path.format(model_init_date,'opxSM_wssfc',1))
 
    # GPM Overpasses
    prs = full_summary(prs, 'Day 0 GPM Overpasses')
    
    # Latest TAFs 
    prs = wxdata_slide(prs, 'Latest (Day 0) TAFs', 'TAF', ['KPAE','KTCM','KHQM'])

    # Summary
    prs = objectives_slide(prs, 'Day 0 Summary')    
    
    # Bumper into day 1 forecast
    prs = bumper_slide(prs, 'Forecast: Day 1', day1_ftime)
    #day1_ftime = model_init_date + timedelta(hours=24)
    #day2_ftime = day2_ftime.replace(hour=0, minute=0, second=0)
      
    # WRF Image --> 500mb Vort
    prs = full_slide_image(prs, 'WRF 500mb Day 1', model_init_date, day1_ftime, link=model_path.format(model_init_date,'opxLG_500mb',5))
    
    # WRF SLP
    prs = full_slide_image(prs, 'WRF SLP Day 1', model_init_date, day1_ftime, link=model_path.format(model_init_date,'opxLG_surface',5))
    
    # WRF Melting level
    prs = full_slide_image(prs, 'WRF Melt. Level Day 1', model_init_date, day1_ftime, link=model_path.format(model_init_date,'opxSM_melt_level',13))
    
    # WRF zoom Precip
    prs = full_slide_image(prs, 'WRF 12hr Prcp (4km) Day 1', model_init_date, day1_ftime, link=model_path.format(model_init_date,'opxSM_precip12hr',13)) 

  
  
    # WRF 3hr zoom Precip
    #prs = full_slide_image(prs, 'WRF 3hr Prcp (4km) Day 1', model_init_date, day1_ftime, link=model_path.format(model_init_date,'opxSM_precip3hr',13)) 
    

    # WRF 4-panel precip
    prs = four_panel_image(prs, 1, model_init_date, link=model_path.format(model_init_date,'opxSM_precip3hr',13))
    # WRF 10m Winds
    prs = full_slide_image(prs, 'WRF 10m Wind (4km) Day 1', model_init_date, day1_ftime, link=model_path.format(model_init_date,'opxSM_wssfc',13))  

    
    # GPM Overpasses
    prs = full_summary(prs, 'Day 1 GPM Overpasses', valid=day1_ftime)
    
    # Timing summary
    #prs = precip_timing_table(prs, 'Day 1 Precip Timing') 
  
    # Possible objectives
    prs = full_summary(prs, 'Day 1 Summary', valid=day1_ftime)

    # Bumper into day 2 forecast
    prs = bumper_slide(prs, 'Forecast: Day 2', day2_ftime)
    #day2_ftime = model_init_date + timedelta(hours=48)

    # WRF Image --> 500mb Vort
    prs = full_slide_image(prs, 'WRF 500mb Day 2', model_init_date, day2_ftime, link=model_path.format(model_init_date,'opxLG_500mb',13))
    
    # WRF SLP
    prs = full_slide_image(prs, 'WRF SLP Day 2', model_init_date, day2_ftime, link=model_path.format(model_init_date,'opxLG_surface',13))
    
    # WRF Melting level
    prs = full_slide_image(prs, 'WRF Melt. Level Day 2', model_init_date, day2_ftime, link=model_path.format(model_init_date,'opxSM_melt_level',37))
    
    # WRF Precip
    prs = full_slide_image(prs, 'WRF 12hr Prcp Day 2', model_init_date, day2_ftime, link=model_path.format(model_init_date,'opxLG_precip12hr',13))    

    # WRF zoom Precip
    #prs = full_slide_image(prs, 'WRF 12hr Prcp (4km) Day 2', model_init_date, day2_ftime, link=model_path.format(model_init_date,'opxSM_precip12hr',37)) 
    
    # WRF 3hr zoom Precip
    #prs = full_slide_image(prs, 'WRF 3hr Prcp (4km) Day 2', model_init_date, day2_ftime, link=model_path.format(model_init_date,'opxSM_precip3hr',37))     
    
    # WRF 4-panel precip
    prs = four_panel_image(prs, 2, model_init_date, link=model_path.format(model_init_date,'opxSM_precip3hr',37)) 
   
    # WRF 10m Winds
    prs = full_slide_image(prs, 'WRF 10m Wind (4km) Day 2', model_init_date, day2_ftime, link=model_path.format(model_init_date,'opxSM_wssfc',37))  

    # GPM Overpasses
    prs = full_summary(prs, 'Day 2 GPM Overpasses', valid=day2_ftime)
    
    # Timing summary
    #prs = precip_timing_table(prs, 'Day 2 Precip Timing')
   
    # Summary
    prs = full_summary(prs, 'Day 2 Summary', valid=day2_ftime)





    # Bumper into day 3+ forecast
    prs = bumper_slide(prs, 'Forecast: Day 3+', day3_ftime)
    
    # WRF Image --> 500mb Vort
    #prs = full_slide_image(prs, 'WRF 500mb Day 3', day2_ftime)
    # WRF 4-panel precip
    prs = four_panel_image(prs, 3, model_init_date, link=model_path.format(model_init_date,'opxSM_precip3hr',61))    
    
    # NAEFS uncertainty   
    prs = full_slide_image(prs, 'NAEFS 500mb and Spread Day 3', model_init_date, link="https://weather.gc.ca/ensemble/naefs/cartes_e.html")
    prs = full_slide_image(prs, 'NAEFS 500mb and Spread Day 4', model_init_date, link="https://weather.gc.ca/ensemble/naefs/cartes_e.html")
    prs = full_slide_image(prs, 'NAEFS 500mb and Spread Day 5', model_init_date, link="https://weather.gc.ca/ensemble/naefs/cartes_e.html")
    # Summary
    prs = full_summary(prs, 'Day 3+ Summary', valid=day3_ftime)

    # Conclusion slide
    prs = full_summary(prs, 'Discussion Summary')

    # Timelines
    prs = full_summary(prs, 'Forecast Timeline')
    prs = full_summary(prs, 'Forecast Timeline')    

    # Insert map slide
    slide_layout = prs.slide_layouts[layout['Blank Slide']]
    slide = prs.slides.add_slide(slide_layout)
    width=9
    left_balanced = (10-width)/2.
    mappath = 'http://www.atmos.washington.edu/~lmadaus/olympex/lynn_ops_domain.png'
    print("Getting ops map")
    try:    
        urllib.request.urlretrieve(mappath, 'lynn_ops_domain.png')
    except AttributeError:
        urllib.urlretrieve(mappath, 'lynn_ops_domain.png')
    except:
        not_found = True
    pic = slide.shapes.add_picture('lynn_ops_domain.png', left=Inches(left_balanced), top=Inches(0.1), width=Inches(width))
    

    # Save the presentation
    prs.save('wxbriefing_{:%Y%m%d%H}.pptx'.format(present_date))

def get_latest_image(product, model_init_date, valid_time=None, within_hours=12):
    """
    Function to get the latest image from a given product,
    but only if it is within within_hours of model_init_date
    
    Returns a tuple with:
    
    (string that is the full address of the image,
    date the image is valid (datetime object))
    
    The product name MUST HAVE an entry in the img_paths dictionary,
    otherwise None is returned
    """
    gfs_hours = {3:60,
                 4:84,
                 5:108}
    naefs_hours = {3:72,
                   4:96,
                   5:120}
    
    if product not in img_paths.keys():
        print("Unable to find path to product:", product)
        return None
    # Parse out the info
    path, ext = img_paths[product]
    origpath = path
    # If this is a path, replace the starttime with the desired time
    # Given as present date (this MUST be a model start time)
    path = path.replace('YYYYMMDDHH',model_init_date.strftime('%Y%m%d%H'))    
    not_found = False
    last_run = False
    # This is a web address
    if product in ['GFS 500mb Day 3', 'NAEFS 500mb and Spread Day 3','NAEFS 500mb and Spread Day 4',\
                'NAEFS 500mb and Spread Day 5']:
        recent_file = ext
        if os.path.exists(recent_file):
            os.system('rm -f {:s}'.format(recent_file))
            
        # Replace the DDDD in the path with the current date
        
        if product.startswith('GFS'):
            fhour = gfs_hours[int(product[-1])]
            fdate = model_init_date.replace(hour=12)
        elif product.startswith('NAEFS'):
            fhour = naefs_hours[int(product[-1])]
            fdate = model_init_date.replace(hour=0)
            
        path = path + '/{:%Y%m%d%H}_{:03d}.gif'.format(fdate,fhour)
        #print(path)
        try:
            urllib.request.urlretrieve(path, recent_file)
        except AttributeError:
            urllib.urlretrieve(path, recent_file)
        except:
            not_found = True            
            
            
        fdate = None
    
    elif ext is None:
        # Just download directly from the path
        fdate = None
        recent_file = path.split('/')[-1]
        if os.path.exists(recent_file):
            os.system('rm -f {:s}'.format(recent_file))
            
        try:    
            urllib.request.urlretrieve(path, recent_file)
        except AttributeError:
            urllib.urlretrieve(path, recent_file)
        except:
            not_found = True
    else:
        if 'WRF' in product:
            fhour = int(path.split('.')[-2][1:])
            fdate = model_init_date + timedelta(hours=fhour)
            pathsplit = path.split('.')
            pathsplit[-3] = fdate.strftime('%Y%m%d%H')
            path = '.'.join(pathsplit)
        else:
            fdate=None
        #print(path)
        if os.path.exists(ext):
            os.system('rm -f {:s}'.format(ext))
                #print(path)
        try:
            urllib.request.urlretrieve(path, ext)
        except AttributeError:
            urllib.urlretrieve(path, ext)
        except:
            not_found = True
        #except:
        #   print("FILE NOT FOUND:")
        #   print(path)
        #   exit(1)
    # Second try here for wrf
    if not_found and "WRF" in product:
        #print("IN HERE!")
        # Reset the paths to the previous model time
        path=origpath
        model_init_date -= timedelta(hours=12)
        path = origpath.replace('YYYYMMDDHH',model_init_date.strftime('%Y%m%d%H'))
        # Subtract 12 hours and see if we can grab that
        fhour = int(path.split('.')[-2][1:])+12
        fdate = model_init_date + timedelta(hours=fhour)
        pathsplit = path.split('.')
        pathsplit[-3] = fdate.strftime('%Y%m%d%H')
        pathsplit[-2] = 'f{:03d}'.format(fhour)
        path = '.'.join(pathsplit)
        
        try:
            urllib.request.urlretrieve(path, ext)
            not_found = False
            last_run = True
        except AttributeError:
            urllib.urlretrieve(path, ext)
            not_found = False
            last_run = True
        except:
            not_found = True
    recent_file = ext

            
    path = ''


        

    # Now check if the file is close in time to presentation date
    """
    if fdate is not None:
        tdiff = model_init_date - fdate
        nhours = tdiff.days * 24 + tdiff.seconds/3600.
        if abs(nhours) > within_hours:
            print("{:s} most recent time is {:%H%MZ %d %b %Y}, skipping".format(product, fdate))
            return None
        print("Found {:s} image valid at {:%H%MZ %d %b %Y}".format(product, fdate))
    else:
    """

    if not_found:
        print("WARNING: Did not find {:s} image".format(product))
        return (None, None)
    else:
        if last_run:
            print("Found {:s} image (PREV RUN WRF)".format(product))
        else:
            print("Found {:s} image".format(product))
        # Return the path
        if path == '':
            return (recent_file, fdate)
        else:
            return ('/'.join((path,recent_file)), fdate)

def add_timeline(slide, curday):
    # Add the timeline to the top of the slide
    shapes = slide.shapes
    # Set colors to use
    daycolor = {0: RGBColor(178, 34, 34),
                1: RGBColor(218, 165, 32),
                2: RGBColor(46, 139, 87),
                3: RGBColor(65, 105, 225)}    
    
    # First, current weather (Day 0)
    left = Inches(0)
    top=Inches(0)
    height = Inches(0.25)
    width = Inches(2.3)
    shape = shapes.add_shape(MSO_SHAPE.PENTAGON,\
        left, top, width, height)
    shape.text = 'Day 0'
    if curday not in [1,2,3,4,5]:
        curday = 0
    elif curday >= 3:
        curday = 3
    
    shape.line.color.rgb = daycolor[0]
    # Now figure out the shading
    shape.fill.solid()
    if curday == 0:
        shape.fill.fore_color.rgb = daycolor[0]
    else:
        shape.fill.fore_color.rgb = RGBColor(200,200,200)
    
    left += width - Inches(0.05)
    width = Inches(2.6)
    # Now for rest of days
    for day in range(1,4):
        
        shape = shapes.add_shape(MSO_SHAPE.CHEVRON,\
            left, top, width, height)
        shape.line.color.rgb = daycolor[day]
        if day == 3:
            shape.text = 'Day {:d}+ ({:%d %b}-)'.format(day, present_date + timedelta(days=day))    
        else:
            shape.text = 'Day {:d} ({:%d %b})'.format(day, present_date + timedelta(days=day))  
        left += width - Inches(0.05)
        shape.fill.solid()
        if curday == day:
            shape.fill.fore_color.rgb = daycolor[curday]
        else:
            shape.fill.fore_color.rgb = RGBColor(200,200,200)  
    
    
    return slide

def precip_timing_table(prs, titletxt):
    """
    Make an empty slide and insert the "timing" chart for front location
    """
    slide_layout = prs.slide_layouts[layout['Title Alone']]
    # Add the slide to the presentation
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = titletxt 
    
    rows = 4
    cols = 4
    width = Inches(9)
    height = Inches(4)
    top = Inches(2)
    left = Inches(0.5)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    table.first_col = True
    table.first_row = True
    # Write text
    table.cell(0,1).text = '<200km Offshore'  
    table.cell(0,2).text = 'At NPOL/Coast'
    table.cell(0,3).text = 'Olympics'
    table.cell(1,0).text = 'Pre-frontal'
    table.cell(2,0).text = 'Frontal Zone'
    table.cell(3,0).text = 'Post-frontal'
 
    # Add timeline
    # First, figure out daynum
    daynum = int(re.search('(\d)', titletxt).groups()[0])
    slide = add_timeline(slide,daynum)    
   
    return prs

def four_panel_image(prs, daynum, model_init_date, link=None):
    # Download all four images
    images = []
    for panel in ['a','b','c','d']:
        results = get_latest_image('WRF 3hr Prcp (4km) Day {:d}{:s}'.format(daynum,panel), model_init_date)
        if results[0] == None:
            images.append([])
        else:
            images.append(results)
    
    # Get a blank slide
    slide_layout = prs.slide_layouts[layout['Title Alone']]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title

    # Now the panels
    lefts = [0,5,0,5]
    tops = [0.15, 0.15, 4, 4]
    text_lefts = [4.2,9.2,4.2,9.2]
    text_tops = [3.0,3.0,6.8,6.8]
    times = ['6Z','12Z','18Z','00Z']
    for pnum, p in enumerate(images):
        if p == []:
            continue
        pic = slide.shapes.add_picture(p[0], left=Inches(lefts[pnum]), top=Inches(tops[pnum]), width=Inches(5))     
        txt = slide.shapes.add_textbox(left=Inches(text_lefts[pnum]), top=Inches(text_tops[pnum]),width=Inches(0.5), height=Inches(0.5))
        tf = txt.text_frame
        para = tf.add_paragraph()
        r = para.add_run()
        r.text = times[pnum]
        r.font.bold = True
        r.font.size = Pt(24)
        if link is not None:
            hlink = r.hyperlink
            hlink.address = link
        
    add_timeline(slide, daynum)
    return prs

def full_slide_image(prs,product,model_init_date, ftime=None, width=None, link=False):
    # Take "product" and make a full-slide image with title out of it
    # Grab the latest image
    results = get_latest_image(product, model_init_date)
    
 
    #imgpath = product
   
    # Get a blank slide layout and add it to the presentation
    slide_layout = prs.slide_layouts[layout['Title Alone']]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    if width is not None:
        title.top=Inches(7.1)
        title.left = Inches(0)
        title.width=Inches(10)
    else:
        title.top=Inches(3)
        title.left = Inches(7.5)
        title.width=Inches(2.0)
    p = title.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = product
    r.font.size=Pt(40)
    if link:
        hlink = r.hyperlink
        hlink.address = link
    if ftime is not None:
        d = p.add_run()
        d.text = '\n\n' + ftime.strftime('%d %b %HZ')
        d.font.size=Pt(28)

    if results[0] is None:
         # Didn't find the image.  Create slide anyway.
         # Add timeline
        dayposs = re.search('Day (\d)', product)
        if dayposs is None:
            curday = 0
        else:
            curday = int(dayposs.groups()[0])
        add_timeline(slide, curday)
        return prs
    imgpath, imgdate = results   
    # Add the image
    if width is not None:
        left_balanced = (10-width)/2.
        pic = slide.shapes.add_picture(imgpath, left=Inches(left_balanced), top=Inches(0.1), width=Inches(width))
    else:
        pic = slide.shapes.add_picture(imgpath, left=Inches(0), top=Inches(0.25), width=Inches(7))
        
    # Add timeline
    dayposs = re.search('Day (\d)', product)
    if dayposs is None:
        curday = 0
    else:
        curday = int(dayposs.groups()[0])
    add_timeline(slide, curday)
    return prs



def bumper_slide(prs, title, date):
    start = date - timedelta(hours=12)
    end = start + timedelta(hours=24)
    # Choose a blank
    slide_layout = prs.slide_layouts[layout['Segue']]
    # Add the slide to the presentation
    slide = prs.slides.add_slide(slide_layout)
    # Change the title
    slide.shapes.title.text = title
    # Subtitle is date
    if 'Current Weather' in title:
        slide.placeholders[1].text = ''
    elif 'Past 24' in title:
        start += timedelta(hours=12)
        end = start + timedelta(hours=24)
        slide.placeholders[1].text = "{:%HZ %d %b %Y} through {:%HZ %d %b %Y}".format(start, end)

    elif 'Day 0' not in title:
        #end = date + timedelta(hours=24)
        slide.placeholders[1].text = "{:%HZ %d %b %Y} through {:%HZ %d %b %Y}".format(start, end)

    else:
        end = date + timedelta(hours=12)
        slide.placeholders[1].text = "Now through {:%HZ %d %b %Y}".format(end)
    
    # Add timeline
    dayposs = re.search('Day (\d)', title)
    if dayposs is None:
        curday = 0
    else:
        curday = int(dayposs.groups()[0])
    add_timeline(slide, curday)


    return prs



def objectives_slide(prs, title):
    # Choose a blank
    slide_layout = prs.slide_layouts[layout['Bullet Slide']]
    # Add the slide to the presentation
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title 
        # Add timeline
    dayposs = re.search('Day (\d)', title)
    if dayposs is None:
        curday = 0
    else:
        curday = int(dayposs.groups()[0])
    add_timeline(slide, curday)
    
    return prs

def full_summary(prs, title, valid=None):
    slide_layout = prs.slide_layouts[layout['Bullet Slide']]
    slide = prs.slides.add_slide(slide_layout)
    if valid is not None:
        title = ' '.join((title, valid.strftime('(%d %b)')))
    slide.shapes.title.text = title
    # Add timeline
    dayposs = re.search('Day (\d)', title)
    if dayposs is None:
        curday = 0
    else:
        curday = int(dayposs.groups()[0])
        
    if title not in ['Discussion Summary', 'Forecast Timeline']:
        add_timeline(slide, curday)
    
    return prs



def get_TAFs(sites):
    main_addr = 'https://www.aviationweather.gov/adds/metars/?station_ids={:s}&std_trans=standard&hoursStr=most+recent+only&chk_tafs=on&submitmet=Submit'
    tafs = []
    for s in sites:
        print("Getting TAF for:", s)
        try:
            import urllib.request
            with urllib.request.urlopen(main_addr.format(s)) as response:
                html = response.read()
        except:
            import urllib2
            response = urllib2.urlopen(main_addr.format(s))
            html = response.read()
        html = str(html)
        part1 = html.split(s.upper())[1]
        part2 = part1.split('</font>')[0]
        datasearch = ' '.join((s.upper(),part2))
        #datasearch = datasearch.replace('\n','\n\r')
        tafs.append(datasearch)
    return tafs

def get_METARs(sites):
    main_addr = 'https://www.aviationweather.gov/adds/metars/?station_ids={:s}&std_trans=standard&chk_metars=on&hoursStr=most+recent+only&submitmet=Submit'
    metars = []
    for s in sites:
        print("Getting METAR for:", s)
        try:
            import urllib.request
            with urllib.request.urlopen(main_addr.format(s)) as response:
                html = response.read()
        except:
            import urllib2
            response = urllib2.urlopen(main_addr.format(s))
            html = response.read()
        html = str(html)
        part1 = html.split(s.upper())[1]
        part2 = part1.split('</FONT>')[0]
        datasearch = ' '.join((s.upper(),part2))
        metars.append(datasearch)
    return metars





def wxdata_slide(prs, title, type='TAF', locs=[]):
    slide_layout = prs.slide_layouts[layout['Bullet Slide']]
    # Add the slide to the presentation
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title 
    # Try to download the requested data
    # If it fails, just return the slide
    # with the title
    try:
        if type == 'TAF':
            data = get_TAFs(locs)
            fs = 14
        elif type == 'METAR':
            data = get_METARs(locs)
            fs = 18
    except:
        return prs
    # Add to the slide
    # Get main text box
    mainbox = slide.placeholders[1]
    tf = mainbox.text_frame
    tf.clear()
    # New paragraph for each station
    for s in data:
        p = tf.add_paragraph()
        p.text = s
        font = p.font
        font.size=Pt(fs)
    # Add timeline
    dayposs = re.search('Day (\d)', title)
    if dayposs is None:
        curday = 0
    else:
        curday = int(dayposs.groups()[0])
    add_timeline(slide, curday)

    return prs


if __name__ == '__main__':
    #get_TAFs(['KPAE','KTCM'])
    build_presentation(model_init_date, present_date)




